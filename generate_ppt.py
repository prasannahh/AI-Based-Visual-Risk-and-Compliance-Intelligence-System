import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette - Premium Dark Theme
    BG_COLOR = RGBColor(15, 23, 42)       # #0F172A Dark Slate
    CARD_BG = RGBColor(30, 41, 59)        # #1E293B Card Background
    CARD_BORDER = RGBColor(51, 65, 85)    # #334155 Card Border
    PRIMARY = RGBColor(14, 165, 233)      # #0EA5E9 Cyan / Light Blue Accent
    SECONDARY = RGBColor(99, 102, 241)    # #6366F1 Indigo Accent
    ACCENT_GREEN = RGBColor(16, 185, 129) # #10B981 Green
    ACCENT_ORANGE = RGBColor(249, 115, 22)# #F97316 Orange
    TEXT_WHITE = RGBColor(248, 250, 252)  # #F8FAFC Main Text
    TEXT_MUTED = RGBColor(148, 163, 184)  # #94A3B8 Secondary Text

    def set_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

    def add_header(slide, title_text, category_text="MILESTONE 2: AI DIGITAL TWIN"):
        # Top banner category tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

        # Divider line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = PRIMARY
        line.line.fill.background()

    def add_card(slide, left, top, width, height, title="", border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        
        if title:
            tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = PRIMARY
        return card

    def add_framed_image(slide, img_path, left, top, width, height):
        # Medium sized framed image box
        if os.path.exists(img_path):
            frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left - 0.05), Inches(top - 0.05), Inches(width + 0.1), Inches(height + 0.1))
            frame.fill.solid()
            frame.fill.fore_color.rgb = CARD_BG
            frame.line.color.rgb = PRIMARY
            frame.line.width = Pt(1.5)
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))

    img_dir = "photos_extracted/projectphoto2"

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)
    
    # Decorative Accent Bar
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.15), Inches(3.2))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = PRIMARY
    accent_bar.line.fill.background()

    # Title Text Frame
    tb = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "MILESTONE 2 PROJECT PRESENTATION"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    p0.space_after = Pt(10)

    p1 = tf.add_paragraph()
    p1.text = "Digital Twin AI - Personal Life Simulation & Decision Assistant"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "AI Core Layer Integration • Predictive ML Engine • Dynamic Multi-Domain Assistants"
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_after = Pt(30)

    p3 = tf.add_paragraph()
    p3.text = "Submitted to: Academic & Project Mentor   |   Team 1   |   Milestone 2 Review"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_GREEN

    # =========================================================================
    # SLIDE 2: Executive Summary & Project Vision
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)
    add_header(slide2, "Executive Summary & Milestone 2 Scope")

    # Card 1: Core Vision
    add_card(slide2, 0.8, 1.6, 5.7, 5.2, "Project Vision: Simulate • Optimize • Evolve")
    tb = slide2.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    bullets1 = [
        ("Digital Twin Concept: ", "Constructs a data-driven personal digital twin modeling user habits, health metrics, financial streams, and study performance."),
        ("Simulate & Predict: ", "Uses Machine Learning algorithms to simulate future outcomes (weight trajectories, savings growth, exam marks, health risks)."),
        ("Actionable Guidance: ", "Delivers personalized, actionable recommendations for diet, workout routines, budgets, and study schedules."),
        ("Unified Decision Support: ", "Acts as an intelligent decision copilot replacing fragmented individual tracking apps.")
    ]
    for i, (b_title, b_desc) in enumerate(bullets1):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + b_title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = PRIMARY
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # Card 2: Milestone 2 Core Achievements
    add_card(slide2, 6.833, 1.6, 5.7, 5.2, "Milestone 2 Key Deliverables")
    tb = slide2.shapes.add_textbox(Inches(7.033), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    bullets2 = [
        ("AI Core Layer Integration: ", "Implemented Machine Learning pipelines using Scikit-Learn, Joblib, and PyTorch for multi-domain prediction."),
        ("Real Dataset Retraining: ", "Trained core risk & performance models on real public datasets (CDC BRFSS 70k+ rows, UCI Obesity, ENSANUT, HF Calorie)."),
        ("4 Integrated AI Assistants: ", "Delivered Health, Financial, Productivity, and Fitness AI Assistants with interactive Streamlit UI."),
        ("Model Registry & Audit: ", "Built automated model versioning, record metrics tracking (AUC, R², MAE), and evaluation logging.")
    ]
    for i, (b_title, b_desc) in enumerate(bullets2):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + b_title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = ACCENT_GREEN
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 3: System Architecture & Tech Stack
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)
    add_header(slide3, "System Architecture & Technology Stack")

    # 4 Architecture Tier Cards
    tiers = [
        ("1. Presentation Tier (UI/UX)", PRIMARY, [
            "Framework: Streamlit Web UI",
            "Responsive Cards & Interactive Tabs",
            "Plotly & Matplotlib Data Visualizations",
            "JWT User Authentication Portal"
        ]),
        ("2. Application & API Tier", SECONDARY, [
            "Framework: FastAPI / Python 3.14",
            "RESTful Endpoint Architecture",
            "JWT Token Session & State Management",
            "Modular Service Controller Layer"
        ]),
        ("3. AI Core & Machine Learning Tier", ACCENT_GREEN, [
            "Frameworks: Scikit-learn, PyTorch, Joblib",
            "Mifflin-St Jeor & BMR Equation Engines",
            "RandomForest, XGBoost, Regression Models",
            "Model Registry (`ai_models/<domain>`)"
        ]),
        ("4. Data & Dataset Pipeline Tier", ACCENT_ORANGE, [
            "Database: SQLite / SQLAlchemy ORM",
            "Public Datasets: CDC BRFSS, UCI, HF",
            "Cached Transformation (`data/real/`)",
            "Synthetic Data Fallback Engine"
        ])
    ]

    for idx, (t_title, color, t_items) in enumerate(tiers):
        col = idx % 2
        row = idx // 2
        left = 0.8 + col * 5.966
        top = 1.6 + row * 2.7
        add_card(slide3, left, top, 5.766, 2.5, t_title, border_color=color)
        
        tb = slide3.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.65), Inches(5.366), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(t_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "✔  " + item
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(4)

    # =========================================================================
    # SLIDE 4: User Authentication & Main Dashboard Overview
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)
    add_header(slide4, "User Authentication & Interactive Main Dashboard")

    # Left: Explanation Box
    add_card(slide4, 0.8, 1.6, 5.7, 5.2, "Dashboard Key Features & Overview")
    tb = slide4.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    dash_features = [
        ("JWT Session Authentication: ", "Secure login/registration system managing user context, session expiration, and personalized profiles."),
        ("Executive KPI Widgets: ", "Displays key user indicators: Age (19), Overall Goal Score (13%), and Active Tracking Days (4)."),
        ("Integrated Multi-Module Cards: ", "Real-time summary widgets showing Savings Trends, Weekly Study Hours breakdown, and Habit Snapshots."),
        ("Behavioral Pattern Engine: ", "Tracks dining out, daily study time, and sleep schedules to detect behavioral shifts and recommend improvements.")
    ]
    for i, (title, desc) in enumerate(dash_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = PRIMARY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # Right: Medium Image (Screenshot 46)
    add_framed_image(slide4, os.path.join(img_dir, "Screenshot (46).png"), 6.833, 1.6, 5.7, 5.2)

    # =========================================================================
    # SLIDE 5: AI Health Assistant - BMI & Calorie Engine
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)
    add_header(slide5, "AI Health Assistant: BMI Calculator & Calorie Engine")

    # Left: Text Explanation
    add_card(slide5, 0.8, 1.6, 5.7, 5.2, "Health Analytics & Caloric Modeling")
    tb = slide5.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    health_features = [
        ("BMI Assessment: ", "Calculates exact BMI (e.g. 24.2 - Healthy Range) from height & weight inputs, delivering tailored nutrition & activity suggestions."),
        ("Mifflin-St Jeor BMR Engine: ", "Computes Basal Metabolic Rate (1,672 kcal) based on biometric parameters."),
        ("Daily Maintenance Requirement: ", "Factors in physical activity levels to determine daily maintenance calories (2,592 kcal)."),
        ("Machine Learning Predictor: ", "Trained on Hugging Face calorie dataset (15,000 records, MAE 1.72 kcal, R² 0.998) to estimate specific burn requirements (2,576 kcal).")
    ]
    for i, (title, desc) in enumerate(health_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = ACCENT_GREEN
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # Right: Medium Images (Screenshot 49 & Screenshot 51)
    add_framed_image(slide5, os.path.join(img_dir, "Screenshot (49).png"), 6.833, 1.6, 5.7, 2.5)
    add_framed_image(slide5, os.path.join(img_dir, "Screenshot (51).png"), 6.833, 4.3, 5.7, 2.5)

    # =========================================================================
    # SLIDE 6: AI Health Assistant - Weight Trajectory & Health Risk
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6)
    add_header(slide6, "AI Health Assistant: Weight Trajectory & Health Risk")

    # Left: Text Explanation
    add_card(slide6, 0.8, 1.6, 5.7, 5.2, "Predictive Weight & Disease Risk Models")
    tb = slide6.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    risk_features = [
        ("Weight Trajectory Forecasting: ", "Scikit-Learn model projects 30-day weight trajectory (70.0 kg to 71.9 kg) based on caloric intake & activity."),
        ("Obesity Risk Classifier: ", "Trained on UCI Obesity dataset (2,110 rows, AUC 1.00, Acc 0.998) -> Evaluates user risk (0% Low Risk)."),
        ("Diabetes Risk Predictor: ", "Trained on CDC BRFSS dataset (70,692 rows, AUC 0.760) -> Predicts diabetic risk (5% Low Risk)."),
        ("Hypertension Risk Classifier: ", "Trained on ENSANUT dataset (3,352 rows, AUC 0.913, Acc 0.927) -> Identifies risk factors (49% Moderate Risk) with preventative tips.")
    ]
    for i, (title, desc) in enumerate(risk_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = PRIMARY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # Right: Medium Images (Screenshot 50 & Screenshot 52)
    add_framed_image(slide6, os.path.join(img_dir, "Screenshot (50).png"), 6.833, 1.6, 5.7, 2.5)
    add_framed_image(slide6, os.path.join(img_dir, "Screenshot (52).png"), 6.833, 4.3, 5.7, 2.5)

    # =========================================================================
    # SLIDE 7: AI Financial Assistant - Savings & Expense Management
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7)
    add_header(slide7, "AI Financial Assistant: Savings & Expense Management")

    # Left: Text Explanation
    add_card(slide7, 0.8, 1.6, 5.7, 5.2, "Financial Analytics & Savings Projections")
    tb = slide7.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    fin_features = [
        ("Financial Transaction Logging: ", "Supports record creation across categories (Salary, Housing, Food, Transport, Investment, Savings)."),
        ("Automated Expense Classifier: ", "ML model (`expense_classifier`) automatically categorizes plain-text descriptions into expense buckets."),
        ("Savings Projection Engine: ", "Forecasts cumulative savings growth over configurable horizons (12-month projection reaching ₹52,000 at ₹4,000/mo rate)."),
        ("Actual vs Projected Analytics: ", "Compares logged actual savings against target trajectories for financial discipline.")
    ]
    for i, (title, desc) in enumerate(fin_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = SECONDARY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # Right: Medium Images (Screenshot 54 & Screenshot 55)
    add_framed_image(slide7, os.path.join(img_dir, "Screenshot (54).png"), 6.833, 1.6, 5.7, 2.5)
    add_framed_image(slide7, os.path.join(img_dir, "Screenshot (55).png"), 6.833, 4.3, 5.7, 2.5)

    # =========================================================================
    # SLIDE 8: AI Financial Assistant - Spending Analysis & Budgeting
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8)
    add_header(slide8, "AI Financial Assistant: Spending Analysis & Budgeting")

    # Left: Text Explanation
    add_card(slide8, 0.8, 1.6, 5.7, 5.2, "Spending Breakdown & Budget Allocation")
    tb = slide8.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    spend_features = [
        ("Categorical Spending Breakdown: ", "Interactive pie chart visualizes expense shares (Total Spent ₹32,000; Top category Housing ₹18,500 / 57.8%)."),
        ("Algorithmic Budget Recommendations: ", "Generates optimal budget limits: Monthly (₹32,000), Weekly (₹7,442), Emergency Fund Target (₹192,000), Savings Goal (₹6,000)."),
        ("Granular Limit Allocation: ", "Establishes category-specific spending caps for Food (8k), Bills (9.6k), Shopping, and Healthcare."),
        ("Balance Forecast Engine: ", "Projects account balance trend over 90 days reaching ₹16,000.")
    ]
    for i, (title, desc) in enumerate(spend_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = PRIMARY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # Right: Medium Images (Screenshot 56 & Screenshot 58)
    add_framed_image(slide8, os.path.join(img_dir, "Screenshot (56).png"), 6.833, 1.6, 5.7, 2.5)
    add_framed_image(slide8, os.path.join(img_dir, "Screenshot (58).png"), 6.833, 4.3, 5.7, 2.5)

    # =========================================================================
    # SLIDE 9: AI Study Assistant - Productivity Tracking & Marks Predictor
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9)
    add_header(slide9, "AI Study Assistant: Productivity & Marks Predictor")

    # Left: Text Explanation
    add_card(slide9, 0.8, 1.6, 5.7, 5.2, "Productivity Tracking & Performance Model")
    tb = slide9.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    study_features = [
        ("Weekly Study Logger: ", "Tracks study hours per subject (Maths, C, English, Biology, CS) with completion rate (89%) and average score (88)."),
        ("Grade & GPA Prediction Model: ", "Trained on UCI Student Performance dataset (1,044 records, R² 0.810, MAE 4.96)."),
        ("Predictive Inputs: ", "Evaluates study hours, days remaining to exam, consistency score, and prior average score."),
        ("Performance Band Output: ", "Predicts expected marks (71.2/100), GPA (3.00), and performance tier ('Good') with curve visualization.")
    ]
    for i, (title, desc) in enumerate(study_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = ACCENT_GREEN
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # Right: Medium Images (Screenshot 61 & Screenshot 62)
    add_framed_image(slide9, os.path.join(img_dir, "Screenshot (61).png"), 6.833, 1.6, 5.7, 2.5)
    add_framed_image(slide9, os.path.join(img_dir, "Screenshot (62).png"), 6.833, 4.3, 5.7, 2.5)

    # =========================================================================
    # SLIDE 10: AI Study Assistant - Weak Subject Detection & Planner
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_background(slide10)
    add_header(slide10, "AI Study Assistant: Weak Subject Detection & Planner")

    # Left: Text Explanation
    add_card(slide10, 0.8, 1.6, 5.7, 5.2, "Weakness Detection & Timetable Generator")
    tb = slide10.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    planner_features = [
        ("Weak Subject Risk Score Ranking: ", "Ranks academic subjects by risk score (e.g. C Programming flagged as highest weakness with Risk Score 0.25, Avg 75/100)."),
        ("Automated Timetable Generator: ", "Generates daily revision schedules balancing target hours (4.0 hrs/day) with upcoming exam dates."),
        ("Time Optimisation Engine: ", "Allocates weekly study hours proportionally to subject weakness rankings."),
        ("Long-Term Performance Trend: ", "Projects 180-day academic score (100) and GPA trajectory (4.0).")
    ]
    for i, (title, desc) in enumerate(planner_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = PRIMARY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # Right: Medium Images (Screenshot 63 & Screenshot 65)
    add_framed_image(slide10, os.path.join(img_dir, "Screenshot (63).png"), 6.833, 1.6, 5.7, 2.5)
    add_framed_image(slide10, os.path.join(img_dir, "Screenshot (65).png"), 6.833, 4.3, 5.7, 2.5)

    # =========================================================================
    # SLIDE 11: AI Fitness Assistant & Habit Tracker
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_background(slide11)
    add_header(slide11, "AI Fitness Assistant & Habit Tracker")

    # Left: Text Explanation
    add_card(slide11, 0.8, 1.6, 5.7, 5.2, "Fitness Scoring & Goal Prediction")
    tb = slide11.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    fit_features = [
        ("Habit Consistency Tracking: ", "Logs habit completion rates (Sleep Schedule 99%, Exercise Frequency 80%, Meal Prep 70%)."),
        ("Composite Fitness Score: ", "Evaluates daily metrics (7,000 steps, 30m exercise, 7h sleep, 2.2L water) producing a composite score (63/100 Intermediate)."),
        ("Personalized Workout Generator: ", "Recommends exercise plans tailored to goals (Intermediate level, 40 min duration, ~350 kcal target, exercises: Jogging, Push-ups, Lunges, Plank hold)."),
        ("Goal Achievement Classifier: ", "Predicts probability of reaching long-term fitness goals (100% High probability to reach score 80 in 90 days).")
    ]
    for i, (title, desc) in enumerate(fit_features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = ACCENT_ORANGE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # Right: Medium Images (Screenshot 69 & Screenshot 70)
    add_framed_image(slide11, os.path.join(img_dir, "Screenshot (69).png"), 6.833, 1.6, 5.7, 2.5)
    add_framed_image(slide11, os.path.join(img_dir, "Screenshot (70).png"), 6.833, 4.3, 5.7, 2.5)

    # =========================================================================
    # SLIDE 12: Machine Learning Models & Real-World Validation
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_background(slide12)
    add_header(slide12, "Machine Learning Models & Real-World Dataset Validation")

    # Left: Model Validation Table Card
    add_card(slide12, 0.8, 1.6, 6.7, 5.2, "Real Dataset Retraining & Hold-Out Validation")
    
    # Table creation
    rows, cols = 6, 4
    left, top, width, height = Inches(1.0), Inches(2.3), Inches(6.3), Inches(4.3)
    table_shape = slide12.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ["Model", "Real Dataset Source", "Sample Size", "Hold-Out Metrics"]
    table_data = [
        ["risk_obesity", "UCI Obesity (Palechor et al.)", "2,110", "AUC 1.00, Acc 0.998"],
        ["risk_diabetes", "CDC BRFSS 2015 Indicators", "70,692", "AUC 0.760, F1 0.705"],
        ["risk_hypertension", "ENSANUT 2021 Survey", "3,352", "AUC 0.913, Acc 0.927"],
        ["performance_pred", "UCI Student Performance", "1,044", "R² 0.810, MAE 4.96"],
        ["workout_calorie", "Hugging Face Calorie Dataset", "15,000", "R² 0.998, MAE 1.72 kcal"]
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = BG_COLOR

    for r, row_data in enumerate(table_data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if r % 2 == 0 else RGBColor(40, 53, 72)
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_WHITE

    # Right: Medium Image (Screenshot 53 / Model Registry Table)
    add_framed_image(slide12, os.path.join(img_dir, "Screenshot (53).png"), 7.8, 1.6, 4.7, 5.2)

    # =========================================================================
    # SLIDE 13: Technical Innovations & Key Engineering Highlights
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_background(slide13)
    add_header(slide13, "Technical Innovations & Engineering Highlights")

    innovations = [
        ("Hybrid Training Architecture", PRIMARY, [
            "Retrained major models on downloaded public datasets (70k+ records).",
            "Maintains synthetic fallback generators (`ai_models/<domain>/synthetic.py`) for schema fields missing in public surveys.",
            "Ensures system reliability even when external datasets have missing features."
        ]),
        ("Modular Domain Pipeline", SECONDARY, [
            "Decoupled ML models into dedicated domains: health, fitness, finance, study.",
            "Joblib serialization for fast inference loading during web UI sessions.",
            "Independent evaluation scripts & automated pytest test suites."
        ]),
        ("Explainable AI & Recommendations", ACCENT_GREEN, [
            "Converts raw probabilistic ML outputs into actionable clinical & financial tips.",
            "Integrates Mifflin-St Jeor biometric equations with ML calorie predictors.",
            "Provides clear disclaimer bounds (health estimates vs medical advice)."
        ]),
        ("Interactive State & Session Management", ACCENT_ORANGE, [
            "Full JWT token validation ensuring secure multi-tenant data isolation.",
            "Reactive Streamlit state synchronization across forms, charts, and tables.",
            "Automated data caching under `data/real/` for fast retraining runs."
        ])
    ]

    for idx, (t_title, color, t_items) in enumerate(innovations):
        col = idx % 2
        row = idx // 2
        left = 0.8 + col * 5.966
        top = 1.6 + row * 2.7
        add_card(slide13, left, top, 5.766, 2.5, t_title, border_color=color)
        
        tb = slide13.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.65), Inches(5.366), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(t_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(4)

    # =========================================================================
    # SLIDE 14: Conclusion & Milestone 3 Roadmap
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_background(slide14)
    add_header(slide14, "Conclusion & Milestone 3 Technical Roadmap")

    # Left: Summary Achievements
    add_card(slide14, 0.8, 1.6, 5.7, 5.2, "Milestone 2 Outcomes Delivered")
    tb = slide14.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    achievements = [
        ("AI Core Fully Functioning: ", "Successfully implemented and validated ML models for Health, Finance, Productivity, and Fitness."),
        ("Real Data Integration: ", "Trained models on real-world datasets achieving up to 1.00 AUC and 0.998 accuracy."),
        ("Polished User Interface: ", "Delivered an interactive web application with dashboards, forecasting charts, and log management."),
        ("Reproducible Pipeline: ", "Automated model retraining script (`ml/train_real.py`) and pytest suite.")
    ]
    for i, (title, desc) in enumerate(achievements):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "✔ " + title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = ACCENT_GREEN
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # Right: Milestone 3 Roadmap
    add_card(slide14, 6.833, 1.6, 5.7, 5.2, "Milestone 3 Planned Roadmap")
    tb = slide14.shapes.add_textbox(Inches(7.033), Inches(2.2), Inches(5.3), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    roadmap = [
        ("Real-World User Pilot: ", "Collect consented real-user feedback to replace synthetic fallback models with empirical data."),
        ("Advanced Deep Learning: ", "Explore Neural Network architectures (PyTorch) for multi-variate time-series forecasting."),
        ("Cloud Deployment: ", "Deploy FastAPI backend & Streamlit web portal to cloud infrastructure (Hugging Face Spaces / AWS)."),
        ("Privacy & Security Audit: ", "Enforce strict HIPAA-compliant data encryption and differential privacy mechanisms.")
    ]
    for i, (title, desc) in enumerate(roadmap):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "➔ " + title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = PRIMARY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_WHITE

    # Save presentation
    output_path = "c:/Users/surya/Desktop/rohith/project/Milestone2_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation successfully saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
